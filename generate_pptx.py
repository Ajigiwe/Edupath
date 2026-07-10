from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
ORANGE = RGBColor(0xC2, 0x41, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
MEDIUM_TEXT = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x10, 0xB9, 0x81)
LIGHT_NAVY = RGBColor(0x2C, 0x4F, 0x7A)
CREAM = RGBColor(0xFF, 0xF7, 0xED)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT,
                bold=False, alignment=PP_ALIGN.LEFT, font_name='Poppins'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16, color=DARK_TEXT,
                          alignment=PP_ALIGN.LEFT, font_name='Inter', bold_first=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_data, tuple):
            text, fs, clr, bld = line_data
        else:
            text = line_data
            fs = font_size
            clr = color
            bld = (bold_first and i == 0)
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(6)
    return tf


def add_shape(slide, left, top, width, height, fill_color, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_accent_bar(slide, left, top, width, height, color=ORANGE):
    return add_shape(slide, left, top, width, height, color)


# ─────────────────────────────────────────────────
# SLIDE 1: TITLE
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, NAVY)

add_accent_bar(slide, 0, 0, 13.333, 0.08, ORANGE)
add_shape(slide, 0, 7.42, 13.333, 0.08, ORANGE)

add_textbox(slide, 1.5, 1.8, 10.3, 1.5, 'EduPath Ghana', font_size=56, color=WHITE, bold=True)
add_textbox(slide, 1.5, 3.3, 10.3, 1.2,
            'Empowering Senior High School Graduates and University\nFreshmen in Ghana to Navigate Admission, Academics, and Career Success.',
            font_size=22, color=RGBColor(0xBB, 0xCC, 0xDD), font_name='Inter')
add_accent_bar(slide, 1.5, 4.7, 3.0, 0.06, ORANGE)
add_textbox(slide, 1.5, 5.2, 10.3, 0.6, 'www.edupathghana.com', font_size=16, color=ORANGE, font_name='Inter')

# ─────────────────────────────────────────────────
# SLIDE 2: THE PROBLEM
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_accent_bar(slide, 0, 0, 13.333, 0.08, ORANGE)

add_textbox(slide, 1.0, 0.4, 11.3, 0.8, 'The Transition Problem in Ghana', font_size=36, color=NAVY, bold=True)
add_shape(slide, 1.0, 1.1, 1.5, 0.06, ORANGE)

# Problem 1
add_textbox(slide, 1.0, 1.8, 5.2, 0.5, 'Information & Admission Deficit', font_size=24, color=ORANGE, bold=True)
add_multiline_textbox(slide, 1.0, 2.4, 5.2, 3.5, [
    ('The Reality', 14, MEDIUM_TEXT, True),
    'Every year, thousands of Ghanaian SHS graduates',
    'face intense rejection from universities like UG',
    'and KNUST simply because they lack clear,',
    'structured information regarding required',
    'aggregate grades and specific cut-off points.',
    '',
    ('Impact', 14, MEDIUM_TEXT, True),
    'Wasted application fees, delayed education,',
    'and misplaced career aspirations.',
], font_size=14, color=DARK_TEXT)

# Divider
add_shape(slide, 6.5, 1.8, 0.04, 4.5, ORANGE)

# Problem 2
add_textbox(slide, 7.0, 1.8, 5.2, 0.5, 'Academic & Resource Scarcity', font_size=24, color=ORANGE, bold=True)
add_multiline_textbox(slide, 7.0, 2.4, 5.2, 3.5, [
    ('The Reality', 14, MEDIUM_TEXT, True),
    'Even after securing entry, enrolled students lack',
    'accessible study materials — syllabi, past questions,',
    'research guides. They are forced to start blind',
    'without adequate resources to learn ahead or',
    'construct clear academic and career plans.',
    '',
    ('Impact', 14, MEDIUM_TEXT, True),
    'High dropout rates, poor academic performance,',
    'and under-prepared graduates entering the workforce.',
], font_size=14, color=DARK_TEXT)

# Bottom stat boxes
for i, (num, label) in enumerate([
    ('400K+', 'WASSCE candidates\nannually'),
    ('60%+', 'Students lack access\nto past questions'),
    ('30%+', 'Freshmen drop out\nin first year'),
]):
    left = 1.0 + i * 4.0
    add_shape(slide, left, 6.0, 3.5, 1.1, LIGHT_GRAY)
    add_textbox(slide, left + 0.3, 6.05, 2.9, 0.5, num, font_size=28, color=ORANGE, bold=True)
    add_textbox(slide, left + 0.3, 6.5, 2.9, 0.5, label, font_size=12, color=MEDIUM_TEXT)

# ─────────────────────────────────────────────────
# SLIDE 3: WHO IS AFFECTED
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_accent_bar(slide, 0, 0, 13.333, 0.08, ORANGE)

add_textbox(slide, 1.0, 0.4, 11.3, 0.8, 'Who Is Affected by This Crisis?', font_size=36, color=NAVY, bold=True)
add_shape(slide, 1.0, 1.1, 1.5, 0.06, ORANGE)

groups = [
    ('WASSCE Candidates', 'Over 400,000 students annually trying to decode competitive aggregates.\n\nExample: KNUST Medicine requires aggregate 6 — most students don\'t know this before applying.', '🎓'),
    ('University Freshmen', 'Students entering higher learning without pre-study resources or foundational course outlines.\n\nThey navigate blindly through their first year with limited guidance.', '📚'),
    ('Anxious Parents', 'Financially burdened guardians funding speculative admissions without strategic predictability.\n\nThey pay fees and WAEC checker costs without knowing their ward\'s real chances.', '👨‍👩‍👧'),
]

for i, (title, desc, icon) in enumerate(groups):
    left = 1.0 + i * 4.0
    add_shape(slide, left, 1.8, 3.5, 4.8, LIGHT_GRAY)
    add_circle(slide, left + 1.1, 2.1, 1.3, ORANGE)
    add_textbox(slide, left + 1.1, 2.1, 1.3, 1.3, icon, font_size=32, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.3, 3.7, 2.9, 0.5, title, font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.3, 4.3, 2.9, 2.0, desc, font_size=12, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────
# SLIDE 4: PLATFORM FEATURES
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_accent_bar(slide, 0, 0, 13.333, 0.08, ORANGE)

add_textbox(slide, 1.0, 0.4, 11.3, 0.8, 'Our Python-Powered Platform Features', font_size=36, color=NAVY, bold=True)
add_shape(slide, 1.0, 1.1, 1.5, 0.06, ORANGE)

features = [
    ('Cut-Off Calculator', 'Instantly calculates admission eligibility using exact school grading equations. No more guessing — students enter grades and see real-time results.'),
    ('Career Pathfinder', 'Maps desired careers back to specific university program alternatives in Ghana, aligned with student performance and interests.'),
    ('Academic Hub', 'Curated repository containing past exams, standard thesis formats, and introductory research guidance for continuous academic support.'),
]

for i, (title, desc) in enumerate(features):
    left = 1.0 + i * 4.0
    # Card
    add_shape(slide, left, 1.8, 3.5, 4.5, LIGHT_GRAY)

    # Number circle
    add_circle(slide, left + 1.3, 2.1, 0.9, ORANGE)
    add_textbox(slide, left + 1.3, 2.1, 0.9, 0.9, str(i + 1), font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, left + 0.3, 3.3, 2.9, 0.7, title, font_size=20, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.3, 4.0, 2.9, 2.0, desc, font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Bottom tagline
add_textbox(slide, 1.0, 6.7, 11.3, 0.5, 'All powered by a smart recommendation engine with WAEC grade mapping and program fit scoring.',
            font_size=14, color=MEDIUM_TEXT, alignment=PP_ALIGN.CENTER, font_name='Inter')

# ─────────────────────────────────────────────────
# SLIDE 5: TARGET CUSTOMERS
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)

add_accent_bar(slide, 0, 0, 13.333, 0.08, ORANGE)

add_textbox(slide, 1.0, 0.4, 11.3, 0.8, 'Identifying Our Core Customers', font_size=36, color=WHITE, bold=True)
add_shape(slide, 1.0, 1.1, 1.5, 0.06, ORANGE)

customers = [
    ('Free Tier Users', 'Secondary school graduates and university\nfreshmen accessing basic admission calculators\nand public forum spaces at no cost.', '#10B981'),
    ('Premium Customers', 'Students and supportive parents who pay\nfor targeted exam mock sets and live career\nmentorship sessions.', '#C2410C'),
    ('B2B School Partners', 'High schools and private colleges subscribing\nto modern aggregate analytics and showcase\nservices for their students.', '#1E3A5F'),
]

for i, (title, desc, accent) in enumerate(customers):
    left = 1.0 + i * 4.0
    color = RGBColor(int(accent[1:3], 16), int(accent[3:5], 16), int(accent[5:7], 16)) if accent.startswith('#') else ORANGE

    add_shape(slide, left, 1.8, 3.5, 4.5, LIGHT_GRAY)
    add_shape(slide, left, 1.8, 3.5, 0.08, color)
    add_textbox(slide, left + 0.3, 2.2, 2.9, 0.6, title, font_size=22, color=color, bold=True)
    add_textbox(slide, left + 0.3, 3.0, 2.9, 2.5, desc, font_size=14, color=DARK_TEXT)

# ─────────────────────────────────────────────────
# SLIDE 6: REVENUE MODEL
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_accent_bar(slide, 0, 0, 13.333, 0.08, ORANGE)

add_textbox(slide, 1.0, 0.4, 11.3, 0.8, 'Our Diversified Revenue Model', font_size=36, color=NAVY, bold=True)
add_shape(slide, 1.0, 1.1, 1.5, 0.06, ORANGE)

revenue_items = [
    ('Subscription Plans', 'Monthly and yearly plans with feature gating.\n4 tiers designed for every segment — Free to Premium Pro (GHS 70/mo).'),
    ('B2B Contracts', 'High schools and private colleges pay annual fees for aggregate analytics, student performance data, and customized dashboards.'),
    ('Exam Mock Sets', 'One-time purchase of targeted mock exam packs for WASSCE and university entrance preparation. Premium add-on for all users.'),
]

for i, (title, desc) in enumerate(revenue_items):
    left = 1.0 + i * 4.0
    add_shape(slide, left, 1.8, 3.5, 3.8, LIGHT_GRAY)

    # Top accent
    add_shape(slide, left, 1.8, 3.5, 0.5, NAVY)
    add_textbox(slide, left + 0.3, 1.9, 2.9, 0.4, title, font_size=20, color=WHITE, bold=True)

    add_textbox(slide, left + 0.3, 2.5, 2.9, 2.5, desc, font_size=13, color=DARK_TEXT)

# Plans summary
plans = [
    ('Free', 'GHS 0', 'Basic access'),
    ('Freemium', 'GHS 15/mo', 'Calculators + forum'),
    ('Premium Basic', 'GHS 40/mo', 'Mock sets + mentorship'),
    ('Premium Pro', 'GHS 70/mo', 'Full suite + priority'),
]

for i, (name, price, desc) in enumerate(plans):
    left = 1.5 + i * 2.8
    add_shape(slide, left, 6.0, 2.5, 1.1, NAVY if i == 3 else (ORANGE if i == 0 else LIGHT_NAVY))
    add_textbox(slide, left + 0.15, 6.05, 2.2, 0.35, name, font_size=14, color=WHITE, bold=True)
    add_textbox(slide, left + 0.15, 6.35, 2.2, 0.3, price, font_size=13, color=RGBColor(0xBB, 0xCC, 0xDD))
    add_textbox(slide, left + 0.15, 6.6, 2.2, 0.3, desc, font_size=10, color=RGBColor(0x99, 0xAA, 0xBB))

# ─────────────────────────────────────────────────
# SLIDE 7: REVENUE PROJECTIONS
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_accent_bar(slide, 0, 0, 13.333, 0.08, ORANGE)

add_textbox(slide, 1.0, 0.4, 11.3, 0.8, 'Projected Revenue Growth', font_size=36, color=NAVY, bold=True)
add_shape(slide, 1.0, 1.1, 1.5, 0.06, ORANGE)

years = [
    ('Year 1', 'Pilot', 'GHS 15,000', 1),
    ('Year 2', 'Launch', 'GHS 45,000', 2),
    ('Year 3', 'Scale', 'GHS 120,000', 3),
    ('Year 4', 'Mature', 'GHS 350,000', 4),
]

bar_max_height = 3.5
bar_base = 6.2
max_value = 350000
bar_width = 2.0
gap = 3.2

for i, (year, phase, amount, level) in enumerate(years):
    left = 1.0 + i * gap
    value = int(amount.replace(',', '').replace('GHS ', ''))
    bar_height = (value / max_value) * bar_max_height

    # Bar
    color = [ORANGE, LIGHT_NAVY, NAVY, GREEN][i]
    add_shape(slide, left + 0.15, bar_base - bar_height, bar_width, bar_height, color)

    # Amount on top
    add_textbox(slide, left, bar_base - bar_height - 0.5, bar_width + 0.3, 0.4, amount, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # Year label
    add_textbox(slide, left, bar_base + 0.15, bar_width + 0.3, 0.35, year, font_size=16, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, bar_base + 0.45, bar_width + 0.3, 0.3, phase, font_size=12, color=MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.0, 6.9, 11.3, 0.4,
            'Expansion in paid subscriptions and nationwide high school B2B contracts will serve as primary scale levers across urban hubs in Ghana.',
            font_size=13, color=MEDIUM_TEXT, alignment=PP_ALIGN.CENTER, font_name='Inter')

# ─────────────────────────────────────────────────
# SLIDE 8: GO-TO-MARKET TIMELINE
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_accent_bar(slide, 0, 0, 13.333, 0.08, ORANGE)

add_textbox(slide, 1.0, 0.4, 11.3, 0.8, 'Our Go-to-Market Timeline', font_size=36, color=NAVY, bold=True)
add_shape(slide, 1.0, 1.1, 1.5, 0.06, ORANGE)

# Timeline line
add_shape(slide, 1.5, 3.55, 10.3, 0.08, ORANGE)

quarters = [
    ('Q1', 'SHS Outreach', 'Conduct informational roadshows across major high schools in Accra and Kumasi. Build awareness through direct student engagement.'),
    ('Q2', 'Social Campaign', 'Partner with micro-influencers and student content creators on TikTok & WhatsApp. Drive viral organic growth.'),
    ('Q3', 'Campus Alliances', 'Partner with University Student Representative Councils (SRCs) during freshers orientation. Onboard thousands at scale.'),
    ('Q4', 'SEO & Blogs', 'Provide comprehensive online aggregate lookup tables to drive steady, high-intent Google search traffic year-round.'),
]

for i, (quarter, title, desc) in enumerate(quarters):
    left = 1.2 + i * 3.0

    # Node circle
    top_y = 2.3 if i % 2 == 0 else 4.0
    add_circle(slide, left + 0.9, 3.15, 0.8, ORANGE if i % 2 == 0 else NAVY)
    add_textbox(slide, left + 0.9, 3.15, 0.8, 0.8, quarter, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Content card
    add_shape(slide, left, top_y, 2.6, 2.5, LIGHT_GRAY)
    add_textbox(slide, left + 0.15, top_y + 0.15, 2.3, 0.4, title, font_size=16, color=NAVY, bold=True)
    add_textbox(slide, left + 0.15, top_y + 0.6, 2.3, 1.7, desc, font_size=11, color=DARK_TEXT)

# ─────────────────────────────────────────────────
# SLIDE 9: CUSTOMER RETENTION
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_accent_bar(slide, 0, 0, 13.333, 0.08, ORANGE)

add_textbox(slide, 1.0, 0.4, 11.3, 0.8, 'Our Customer Retention Mechanisms', font_size=36, color=NAVY, bold=True)
add_shape(slide, 1.0, 1.1, 1.5, 0.06, ORANGE)

retention = [
    ('Gamified Coursework Paths', 'Keep students motivated by tracking pre-university course completion with digital credential badges and progress milestones.'),
    ('Interactive Community Forums', 'Tailor discussions based on target universities (e.g., KNUST Engineering or UG Legon Business) to maintain active peer dialogue.'),
    ('Continuous Alumni Mentoring', 'Connect junior students with experienced seniors on campus for real-time tips, syllabi review, and tutoring support.'),
    ('Dynamic Resource Updating', 'Constantly refresh past questions and project guides, expanding user lifespan from entry up through graduation.'),
]

for i, (title, desc) in enumerate(retention):
    row = i // 2
    col = i % 2
    left = 1.0 + col * 5.8
    top = 1.6 + row * 2.6

    add_shape(slide, left, top, 5.5, 2.2, LIGHT_GRAY)
    # Left accent bar
    add_shape(slide, left, top, 0.08, 2.2, ORANGE)
    add_textbox(slide, left + 0.4, top + 0.25, 4.8, 0.4, title, font_size=18, color=NAVY, bold=True)
    add_textbox(slide, left + 0.4, top + 0.75, 4.8, 1.2, desc, font_size=13, color=DARK_TEXT)

# ─────────────────────────────────────────────────
# SLIDE 10: Q&A / THANK YOU
# ─────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)

add_accent_bar(slide, 0, 0, 13.333, 0.08, ORANGE)
add_shape(slide, 0, 7.42, 13.333, 0.08, ORANGE)

add_textbox(slide, 1.5, 1.8, 10.3, 1.2, 'Questions & Answers', font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, 5.5, 3.1, 2.3, 0.06, ORANGE)

add_textbox(slide, 1.5, 3.5, 10.3, 0.8,
            'Join us in guiding the next generation of Ghanaian scholars\ntowards university and career milestones.',
            font_size=20, color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER, font_name='Inter')

add_textbox(slide, 1.5, 5.0, 10.3, 0.6, 'www.edupathghana.com', font_size=20, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 6.0, 10.3, 0.5, 'Thank You', font_size=18, color=RGBColor(0x99, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER, font_name='Inter')

# ─────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────
output_path = r'C:\Users\ABCD\Desktop\atio\educate460\EduPath_Ghana_Pitch_Deck.pptx'
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
